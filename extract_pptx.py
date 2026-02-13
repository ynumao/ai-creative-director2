from pptx import Presentation
import sys

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text_content = []
    for i, slide in enumerate(prs.slides):
        text_content.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    text_content.append(" | ".join(cell.text_frame.text for cell in row.cells))
    return "\n".join(text_content)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py <path_to_pptx>")
    else:
        print(extract_text_from_pptx(sys.argv[1]))
